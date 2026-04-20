import io
import os
import base64
import requests
from fastapi import APIRouter, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import Optional
from pptx import Presentation
from pptx.oxml.ns import qn

router = APIRouter()

NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

TEMPLATE_PATH = os.path.join(
    os.path.dirname(__file__), "..", "templates", "base-presentation.pptx"
)

# Chanel X Booth — utilisé comme BASE TEMPLATE (layout cabines/kiosk/goodies correct)
CHANEL_PATH = os.path.join(
    os.path.dirname(__file__), "..", "..", "exemple", "Chanel X Booth.pptx"
)


class GenerateRequest(BaseModel):
    brand_name: str
    website: Optional[str] = None
    primary_color: Optional[str] = None    # hex sans '#', ex: "C5A028"
    secondary_color: Optional[str] = None  # hex sans '#', ex: "F5F0E8"
    logo_url: Optional[str] = None         # URL publique ou data:image/... du logo (wordmark)
    logo_icon_url: Optional[str] = None     # URL de l'icône/symbole (ex: le «N» de Netflix)
    # cover_style: style de la page de couverture
    #   "brand"   → fond plein couleur primaire + logo centré
    #   "split"   → bandeau vertical primaire gauche + fond secondaire droite
    #   "minimal" → fond clair secondaire + barre primaire en bas
    cover_style: Optional[str] = "brand"
    # strip_style: style du contour des bandes photo (Freeform 5 & 7, slide 3)
    #   "primary"   → bordure épaisse couleur primaire
    #   "secondary" → bordure fine couleur secondaire
    #   "none"      → aucune bordure (défaut Chanel)
    strip_style: Optional[str] = "none"
    cabine_top_url: Optional[str] = None       # legacy
    cabine_bottom_url: Optional[str] = None    # legacy
    cabine_ronde_url: Optional[str] = None     # Gemini image-to-image — Freeform 25
    cabine_carree_url: Optional[str] = None    # Gemini image-to-image — Freeform 24
    kiosk_url: Optional[str] = None
    goodies_top_url: Optional[str] = None
    goodies_bottom_url: Optional[str] = None
    # 3 images brand pour les slots de la cabine carrée (optionnel — générées si absentes)
    cabine_slot_urls: Optional[list] = None


def get_image_bytes(url_or_data: str) -> bytes:
    """Get image bytes from either a base64 data URL or an http URL."""
    if url_or_data.startswith("data:"):
        # data:image/png;base64,<b64>
        header, b64 = url_or_data.split(",", 1)
        return base64.b64decode(b64)
    return download_image(url_or_data)


def download_image(url: str, retries: int = 5) -> bytes:
    """Download image with retry/backoff — handles Pollinations 429/500."""
    import time
    for attempt in range(retries + 1):
        try:
            resp = requests.get(url, timeout=60)
            if resp.status_code == 200:
                return resp.content
            if resp.status_code in (429, 500, 503) and attempt < retries:
                delay = min(2 ** attempt * 2, 20)  # 2s, 4s, 8s, 16s, 20s
                print(f"Pollinations {resp.status_code} — retry {attempt+1}/{retries} in {delay}s")
                time.sleep(delay)
                continue
            resp.raise_for_status()
        except requests.exceptions.Timeout:
            if attempt < retries:
                print(f"Pollinations timeout — retry {attempt+1}/{retries}")
                time.sleep(5)
                continue
            raise
    raise RuntimeError(f"download_image failed after {retries} retries: {url}")


def get_shape(slide, name: str):
    """Return shape by name, searching inside groups too."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
        if shape.shape_type == 6:  # GROUP
            for child in shape.shapes:
                if child.name == name:
                    return child
    return None


def replace_blip(slide_part, shape, img_bytes: bytes) -> bool:
    """Replace the blipFill image inside a freeform or group shape."""
    spTree = shape._element
    blip = spTree.find(".//" + qn("a:blip"))
    if blip is None:
        return False
    rid = blip.get("{%s}embed" % NS_R)
    if not rid:
        return False
    try:
        slide_part._rels[rid].target_part._blob = img_bytes
        return True
    except Exception:
        return False


def replace_blip_contained(slide_part, shape, img_bytes: bytes) -> bool:
    """Replace blipFill image AND adjust fillRect so the image is letterboxed/pillarboxed
    (contained, not stretched) within the shape bounds. Requires Pillow."""
    # First do the regular blip swap
    if not replace_blip(slide_part, shape, img_bytes):
        return False

    # Read image dimensions to compute aspect-ratio-correct fillRect padding
    try:
        from PIL import Image as _PILImage
        import io as _io
        img = _PILImage.open(_io.BytesIO(img_bytes))
        iw, ih = img.size
        img.close()
    except Exception:
        return True  # blip replaced, just skip padding

    sw = shape.width
    sh = shape.height
    if sw <= 0 or sh <= 0 or iw <= 0 or ih <= 0:
        return True

    shape_ar = sw / sh
    img_ar   = iw / ih

    # fillRect values are in thousandths of a percent (100000 = 100%)
    # Positive = inset (shrink), negative = expand
    if img_ar > shape_ar:
        # Image is wider than shape → fit by width, add top/bottom padding
        rendered_frac = shape_ar / img_ar          # fraction of shape height used
        pad = int((1 - rendered_frac) / 2 * 100000)
        fr = {"l": "0", "t": str(pad), "r": "0", "b": str(pad)}
    else:
        # Image is taller than shape → fit by height, add left/right padding
        rendered_frac = img_ar / shape_ar
        pad = int((1 - rendered_frac) / 2 * 100000)
        fr = {"l": str(pad), "t": "0", "r": str(pad), "b": "0"}

    # Update the fillRect element in blipFill
    fill_rect_el = shape._element.find(".//" + qn("a:fillRect"))
    if fill_rect_el is not None:
        for k, v in fr.items():
            fill_rect_el.set(k, v)

    return True


def replace_text(slide, old: str, new: str):
    """Replace text occurrences across all shapes in a slide, including group children."""
    def _replace_in_shape(s):
        if s.has_text_frame:
            for para in s.text_frame.paragraphs:
                for run in para.runs:
                    if old in run.text:
                        run.text = run.text.replace(old, new)
        if s.shape_type == 6:  # GROUP — recurse into children
            for child in s.shapes:
                _replace_in_shape(child)

    for shape in slide.shapes:
        _replace_in_shape(shape)


def _delete_slide(prs, slide_idx: int):
    """Remove a slide from the presentation by index."""
    # prs._element is the CT_Presentation XML element in python-pptx
    sld_id_lst = prs._element.find(qn("p:sldIdLst"))
    sld_ids = sld_id_lst.findall(qn("p:sldId"))
    sld_id_lst.remove(sld_ids[slide_idx])


def _rebrand_tarifs_slide(prs, slide_idx: int, primary_hex: str, logo_bytes: Optional[bytes]):
    """
    Slide "Nos tarifs" : adapte fond + watermark logo à la marque.
    - Fond : couleur primaire de la marque (au lieu du noir Chanel)
    - Freeform 4 : grand watermark logo semi-transparent (alpha 29%) → remplacé par le logo brand
    - Freeform 5 : petit logo en bas à droite → remplacé par le logo brand
    """
    if slide_idx >= len(prs.slides):
        return
    slide = prs.slides[slide_idx]

    # 1. Changer la couleur de fond
    try:
        bg_pr = slide._element.find('.//' + qn('p:bgPr'))
        if bg_pr is not None:
            solid = bg_pr.find(qn('a:solidFill'))
            if solid is not None:
                srgb = solid.find(qn('a:srgbClr'))
                if srgb is not None:
                    srgb.set('val', primary_hex.lstrip('#').upper())
    except Exception as e:
        print(f'Warning: tarifs bg color change failed: {e}')

    if not logo_bytes:
        return

    # 2. Remplacer le blip dans Freeform 4 (watermark) et Freeform 5 (logo icon)
    # Même technique que replace_blip : on écrase directement le blob de la relation.
    for shape_name in ('Freeform 4', 'Freeform 5'):
        shape = get_shape(slide, shape_name)
        if shape is None:
            continue
        try:
            sp = shape._element
            blip = sp.find('.//' + qn('a:blip'))
            if blip is None:
                continue
            rid = blip.get('{%s}embed' % NS_R)
            if rid:
                slide.part._rels[rid].target_part._blob = logo_bytes
            # Supprimer les extensions SVG qui ne s'appliquent pas au nouveau logo
            ext_lst = blip.find(qn('a:extLst'))
            if ext_lst is not None:
                blip.remove(ext_lst)
        except Exception as e:
            print(f'Warning: tarifs logo replace failed for {shape_name}: {e}')


def _set_shape_solid_fill(shape, hex_color: str):
    """Replace any fill on a shape with a solid color (hex without '#')."""
    from lxml import etree as _etree
    sp_pr = shape._element.find(qn("p:spPr"))
    if sp_pr is None:
        return
    for tag in ("a:blipFill", "a:solidFill", "a:gradFill", "a:pattFill", "a:noFill"):
        el = sp_pr.find(qn(tag))
        if el is not None:
            sp_pr.remove(el)
    solid = _etree.SubElement(sp_pr, qn("a:solidFill"))
    srgb = _etree.SubElement(solid, qn("a:srgbClr"))
    srgb.set("val", hex_color.lstrip("#").upper())


def inject_picture_at_shape(slide, shape_name: str, img_bytes: bytes):
    """Add an image as a new picture at the exact position/size of an existing shape."""
    import io
    target = get_shape(slide, shape_name)
    if target is None:
        print(f"Warning: shape '{shape_name}' not found for picture injection")
        return
    try:
        slide.shapes.add_picture(
            io.BytesIO(img_bytes),
            target.left, target.top,
            target.width, target.height,
        )
    except Exception as e:
        print(f"Warning: inject_picture_at_shape failed for '{shape_name}': {e}")


def set_cover_overlay_color(slide, hex_color: str):
    """Change the solidFill of Freeform 3 (inside Group 2) on the cover slide to the brand primary color."""
    clean = hex_color.lstrip("#").upper()
    if len(clean) != 6:
        return
    for shape in slide.shapes:
        if shape.name == "Group 2" and shape.shape_type == 6:
            for child in shape.shapes:
                if child.name == "Freeform 3":
                    solid = child._element.find(".//" + qn("a:solidFill"))
                    if solid is not None:
                        srgb = solid.find(qn("a:srgbClr"))
                        if srgb is not None:
                            srgb.set("val", clean)
                    return


def update_cover_texts(slide):
    """Update cover slide textbox labels to match the Booth template standard."""
    # Chanel cover uses TextBox 4 and 5 ; neutral uses TextBox 6 and 7
    mapping = {
        "TextBox 4": "PHOTOBOOTH | BRANDED BOOTH",
        "TextBox 5": "PRÉSENTATION DES SERVICES | LOCATION LONGUE DURÉE 2025/2026",
        "TextBox 6": "PHOTOBOOTH | BRANDED BOOTH",
        "TextBox 7": "PRÉSENTATION DES SERVICES | LOCATION LONGUE DURÉE 2025/2026",
    }
    for shape in slide.shapes:
        new_text = mapping.get(shape.name)
        if new_text and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                runs = para.runs
                if runs:
                    runs[0].text = new_text
                    for run in runs[1:]:
                        run.text = ""
                    new_text = ""  # only update first paragraph


# ── Cover helpers ──────────────────────────────────────────────────────────

def _hex(raw: Optional[str], fallback: str = "0D0D0D") -> str:
    """Normalize a hex color string (strip '#', uppercase, pad to 6)."""
    if not raw:
        return fallback.upper()
    clean = raw.lstrip("#").upper()
    return clean if len(clean) == 6 else fallback.upper()


def _is_dark(hex6: str) -> bool:
    """Return True if the color is dark (luminance < 0.4)."""
    try:
        r, g, b = int(hex6[0:2], 16), int(hex6[2:4], 16), int(hex6[4:6], 16)
        lum = (0.299 * r + 0.587 * g + 0.114 * b) / 255
        return lum < 0.55
    except Exception:
        return True


def _set_slide_bg(slide, hex6: str) -> None:
    """Fill the slide background with a solid color."""
    from pptx.dml.color import RGBColor
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex6)


def _add_rect(slide, left, top, width, height, hex6: str, line=False) -> None:
    """Add a filled rectangle with no border (or thin border if line=True)."""
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(hex6)
    if line:
        shape.line.color.rgb = RGBColor.from_string(hex6)
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()


def _add_textbox(slide, left, top, width, height,
                 text: str, hex6_color: str, font_size_pt: float,
                 bold=False, letter_spacing_pt=1.5) -> None:
    """Add a simple single-line text box."""
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(hex6_color)
    # Approximate letter spacing via character spacing XML
    try:
        from lxml import etree as _et
        from pptx.oxml.ns import qn as _qn
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(letter_spacing_pt * 100)))
    except Exception:
        pass


def _add_logo(slide, logo_bytes: bytes, cx: int, cy: int,
              max_w: int, max_h: int) -> None:
    """Add logo centered at (cx, cy) within max_w × max_h bounds."""
    lw, lh = max_w, max_h
    try:
        from PIL import Image as _PILImage
        img = _PILImage.open(io.BytesIO(logo_bytes))
        iw, ih = img.size
        img.close()
        scale = min(max_w / iw, max_h / ih)
        lw, lh = int(iw * scale), int(ih * scale)
    except Exception:
        pass
    left = cx - lw // 2
    top  = cy - lh // 2
    slide.shapes.add_picture(io.BytesIO(logo_bytes), left, top, lw, lh)


def _recolor_logo(logo_bytes: bytes, bg_hex: str) -> bytes:
    """
    Recolorie le logo pour qu'il soit visible sur bg_hex (la couleur de fond).
    Si le contraste est insuffisant (ratio < 2.5) et que le logo est monochrome,
    recolorie vers blanc (fond sombre) ou noir (fond clair).
    Les logos multicolores (Desigual, etc.) sont conservés intacts.
    """
    try:
        from PIL import Image as _PIL
        img = _PIL.open(io.BytesIO(logo_bytes)).convert("RGBA")
        data = list(img.getdata())

        # Pixels opaques uniquement
        opaque = [(r, g, b) for (r, g, b, a) in data if a > 30]
        if not opaque:
            return logo_bytes

        # Luminance moyenne du logo
        def _lum(r, g, b):
            return (0.299 * r + 0.587 * g + 0.114 * b) / 255

        logo_lum = sum(_lum(*p) for p in opaque) / len(opaque)

        # Luminance du fond réel sur lequel le logo sera posé
        clean_bg = bg_hex.lstrip("#").upper().zfill(6)[:6]
        bg_lum = (
            0.299 * int(clean_bg[0:2], 16)
            + 0.587 * int(clean_bg[2:4], 16)
            + 0.114 * int(clean_bg[4:6], 16)
        ) / 255

        # Ratio de contraste (formule WCAG simplifiée)
        lighter = max(logo_lum, bg_lum) + 0.05
        darker  = min(logo_lum, bg_lum) + 0.05
        contrast = lighter / darker

        # Si contraste suffisant (≥ 3.0) → garder le logo original
        if contrast >= 3.0:
            return logo_bytes

        # Contraste insuffisant → vérifier si le logo est multicolore
        import colorsys
        hues = [colorsys.rgb_to_hsv(r/255, g/255, b/255)[0]
                for (r, g, b) in opaque
                if max(r, g, b) - min(r, g, b) > 20]  # pixels saturés seulement

        is_multicolor = len(hues) > 50 and (max(hues) - min(hues)) > 0.15

        # Logo multicolore avec mauvais contraste → on ne recolorie pas
        if is_multicolor:
            return logo_bytes

        # Logo monochrome/neutre avec mauvais contraste → recolorie vers couleur contrastante
        # Fond sombre → logo blanc ; fond clair → logo noir
        if bg_lum < 0.5:
            r_t, g_t, b_t = 255, 255, 255   # blanc sur fond sombre
        else:
            r_t, g_t, b_t = 17,  17,  17    # quasi-noir sur fond clair
        new_data = [
            (r_t, g_t, b_t, a) if a > 30 else (r, g, b, 0)
            for (r, g, b, a) in data
        ]
        img.putdata(new_data)
        out = io.BytesIO()
        img.save(out, format="PNG")
        return out.getvalue()
    except Exception:
        return logo_bytes


def _hide_existing_cover_shapes(slide) -> None:
    """
    Masque les formes originales du slide cover Chanel en les rendant
    invisibles (transparentes) pour ne garder que nos nouvelles formes.
    """
    from lxml import etree as _etree
    for shape in slide.shapes:
        if shape.name in ("Freeform 2", "Freeform 3"):
            try:
                _set_shape_solid_fill(shape, "000000")  # on le peindra ensuite
                sp_pr = shape._element.find(qn("p:spPr"))
                if sp_pr is not None:
                    # Remplacer par noFill
                    for tag in ("a:blipFill", "a:solidFill", "a:gradFill",
                                "a:pattFill", "a:noFill"):
                        el = sp_pr.find(qn(tag))
                        if el is not None:
                            sp_pr.remove(el)
                    _etree.SubElement(sp_pr, qn("a:noFill"))
            except Exception:
                pass
        # TextBox 4 (bas-droite) et TextBox 5 (bas-gauche) sont conservés tels quels :
        # leur police/espacement du template Chanel est préservé.
        # Seule leur couleur sera mise à jour via _recolor_cover_texts().


def _rebrand_photo_strips(
    slide, primary: str, secondary: str,
    brand_name: str, logo_bytes: Optional[bytes],
    logo_icon_bytes: Optional[bytes] = None
) -> None:
    """
    Rebrande les strips photo du slide Photobooth Classique (slide 3) :
    - Freeform 5 (strip vertical, 1 grande photo) :
        remplace la barre noire du bas (~y=410→bas) avec la couleur primaire
        de la marque + le nom de la marque en texte centeré.
    - Freeform 7 (strip 4 photos en grille) :
        remplace la zone blanche du bas (~y=430→bas) avec la couleur
        secondaire (ou blanc cassé) + le logo de la marque centeré.
    Tout le reste (photos, bords arrondis) est conservé intact.
    """
    from PIL import Image as _PIL, ImageDraw
    import os

    def _get_img_and_part(shape):
        """Retourne (blob bytes, image_part) à partir du blip de la shape."""
        sp = shape._element
        blip = sp.find('.//' + qn('a:blip'))
        if blip is None:
            return None, None
        NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        rId = blip.get(f'{{{NS}}}embed')
        if not rId:
            return None, None
        for rel in slide.part.rels.values():
            if rel.rId == rId and rel.reltype.endswith('/image'):
                return rel._target.blob, rel._target
        return None, None

    def _iter(shapes):
        for s in shapes:
            yield s
            if s.shape_type == 6:
                yield from _iter(s.shapes)

    # Couleurs normalisées
    # Freeform 5 : fond PRIMAIRE + logo original
    # Freeform 7 : fond SECONDAIRE clair + logo original (le vrai logo Brandfetch)
    pr      = _hex(primary,   "111111")
    sc_orig = _hex(secondary, "F5F3EE")  # secondaire avant forçage
    sc      = sc_orig if not _is_dark(sc_orig) else "F5F3EE"   # fond clair garanti
    pr_rgb  = tuple(int(pr[i:i+2], 16) for i in (0, 2, 4))
    sc_rgb  = tuple(int(sc[i:i+2], 16) for i in (0, 2, 4))

    for shape in _iter(slide.shapes):
        if shape.name not in ("Freeform 5", "Freeform 7"):
            continue
        blob, img_part = _get_img_and_part(shape)
        if not blob or not img_part:
            continue

        img = _PIL.open(io.BytesIO(blob)).convert("RGBA")
        w, h = img.size
        draw = ImageDraw.Draw(img)

        if shape.name == "Freeform 5":
            # ── Barre bas : fond PRIMAIRE + logo recolorisé pour contraste ──
            BAR_Y = int(h * 0.762)  # ~410/538
            draw.rectangle([(0, BAR_Y), (w, h)], fill=pr_rgb + (255,))
            zone_h = h - BAR_Y

            if logo_bytes:
                try:
                    # Logo recolorisé pour être visible sur fond primaire
                    _logo_vis = _recolor_logo(logo_bytes, pr)
                    logo_img = _PIL.open(io.BytesIO(_logo_vis)).convert("RGBA")
                    max_lw = int(w * 0.52)
                    max_lh = int(zone_h * 0.65)
                    lw, lh = logo_img.size
                    scale  = min(max_lw / lw, max_lh / lh)
                    new_lw = max(1, int(lw * scale))
                    new_lh = max(1, int(lh * scale))
                    logo_img = logo_img.resize((new_lw, new_lh), _PIL.LANCZOS)
                    lx = (w - new_lw) // 2
                    ly = BAR_Y + (zone_h - new_lh) // 2
                    img.paste(logo_img, (lx, ly), logo_img)
                except Exception as e:
                    print(f"Warning: logo paste on Freeform 5 failed: {e}")

        elif shape.name == "Freeform 7":
            # ── Zone logo bas : fond SECONDAIRE clair + icône/symbole Brandfetch ──
            LOGO_Y = int(h * 0.799)  # ~430/538
            draw.rectangle([(0, LOGO_Y), (w, h)], fill=sc_rgb + (255,))
            zone_h = h - LOGO_Y

            # Utilise l'icône (ex: «N» Netflix) si dispo, sinon logo principal
            _icon = logo_icon_bytes if logo_icon_bytes else logo_bytes
            if _icon:
                try:
                    logo_img = _PIL.open(io.BytesIO(_icon)).convert("RGBA")
                    max_lw = int(w * 0.70)   # plus large que Freeform 5 car c'est le wordmark
                    max_lh = int(zone_h * 0.60)
                    lw, lh = logo_img.size
                    scale  = min(max_lw / lw, max_lh / lh)
                    new_lw = max(1, int(lw * scale))
                    new_lh = max(1, int(lh * scale))
                    logo_img = logo_img.resize((new_lw, new_lh), _PIL.LANCZOS)
                    lx = (w - new_lw) // 2
                    ly = LOGO_Y + (zone_h - new_lh) // 2
                    img.paste(logo_img, (lx, ly), logo_img)
                except Exception as e:
                    print(f"Warning: logo paste on Freeform 7 failed: {e}")

        # Réinjection du blob modifié
        out = io.BytesIO()
        img.save(out, format="PNG")
        img_part._blob = out.getvalue()


def _style_photo_strips(slide, primary: str, secondary: str, strip_style: str) -> None:
    """
    Applique un contour coloré sur les Freeform 5 et Freeform 7 (slide 3 —
    les bandes photo). Les bords arrondis et le contenu image sont conservés.
      "primary"   → trait 3 pt couleur primaire de la marque
      "secondary" → trait 1.5 pt couleur secondaire
      "none"      → aucun trait (transparent)
    """
    from lxml import etree as _et
    from pptx.oxml.ns import qn as _qn

    STYLE = (strip_style or "none").lower().strip()
    TARGET_NAMES = {"Freeform 5", "Freeform 7"}

    # Récupére tous les shapes y compris enfants de groupes
    def iter_shapes(shapes):
        for s in shapes:
            yield s
            if s.shape_type == 6:  # GROUP
                yield from iter_shapes(s.shapes)

    for shape in iter_shapes(slide.shapes):
        if shape.name not in TARGET_NAMES:
            continue
        sp_pr = shape._element.find(_qn("p:spPr"))
        if sp_pr is None:
            continue
        # Supprimer l'ancien a:ln s'il existe
        old_ln = sp_pr.find(_qn("a:ln"))
        if old_ln is not None:
            sp_pr.remove(old_ln)

        if STYLE == "none":
            # Explicit noFill border so PowerPoint doesn’t inherit theme line
            ln = _et.SubElement(sp_pr, _qn("a:ln"))
            _et.SubElement(ln, _qn("a:noFill"))
        else:
            color_hex = primary if STYLE == "primary" else secondary
            # Normalise: 6 hex chars, no '#'
            color_hex = color_hex.lstrip("#").upper().zfill(6)[:6]
            width_emu = 38100 if STYLE == "primary" else 19050  # 3 pt / 1.5 pt
            ln = _et.SubElement(sp_pr, _qn("a:ln"), w=str(width_emu))
            solid = _et.SubElement(ln, _qn("a:solidFill"))
            _et.SubElement(solid, _qn("a:srgbClr"), val=color_hex)


def _bring_cover_texts_to_front(slide) -> None:
    """
    Déplace TextBox 4 et TextBox 5 à la fin du spTree (z-order le plus haut)
    pour qu'ils apparaissent au-dessus des rectangles ajoutés dynamiquement.
    """
    sp_tree = slide.shapes._spTree
    tb4_xml = tb5_xml = None
    for shape in slide.shapes:
        if shape.name == "TextBox 4":
            tb4_xml = shape._element
        elif shape.name == "TextBox 5":
            tb5_xml = shape._element
    for xml_el in [tb4_xml, tb5_xml]:
        if xml_el is not None and xml_el in sp_tree:
            sp_tree.remove(xml_el)
            sp_tree.append(xml_el)


def _recolor_cover_texts(slide, color_tb5: str, color_tb4: str) -> None:
    """
    Recolorie les runs de TextBox 5 (bas-gauche : PRESENTATION DES SERVICES…)
    et TextBox 4 (bas-droite : PHOTOBOOTH | BRANDED BOOTH).
    La police, la taille, l'espacement et le layout du template Chanel sont préservés.
    """
    from pptx.dml.color import RGBColor
    targets = {"TextBox 4": color_tb4, "TextBox 5": color_tb5}
    for shape in slide.shapes:
        color = targets.get(shape.name)
        if color and hasattr(shape, "text_frame"):
            try:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = RGBColor.from_string(color)
            except Exception:
                pass


# ── 3 styles de cover ──────────────────────────────────────────────────────

def _cover_brand(slide, W, H, primary: str, secondary: str,
                 logo_bytes: Optional[bytes], brand_name: str) -> None:
    """
    Style A — "brand" : fond plein couleur primaire.
    Logo centré. Ligne fine + sous-titre en bas à gauche.
    """
    on_dark = _is_dark(primary)
    text_color = "FFFFFF" if on_dark else "111111"
    accent     = secondary if not _is_dark(secondary) else "FFFFFF"

    # Fond plein
    _set_slide_bg(slide, primary)

    # Ligne horizontale fine sous le logo (80 % largeur)
    line_w = int(W * 0.80)
    line_h = int(H * 0.003)
    _add_rect(slide, (W - line_w) // 2, int(H * 0.64), line_w, line_h, text_color)

    # Logo centré — recolorisé pour contraster avec le fond (on passe la couleur du fond)
    if logo_bytes:
        logo_colored = _recolor_logo(logo_bytes, primary)
        _add_logo(slide, logo_colored,
                  cx=W // 2, cy=int(H * 0.42),
                  max_w=int(W * 0.32), max_h=int(H * 0.35))

    # Texte haut de marque — discret, espacé
    _add_textbox(slide,
                 left=int(W * 0.08), top=int(H * 0.06),
                 width=int(W * 0.84), height=int(H * 0.08),
                 text=brand_name.upper(),
                 hex6_color=text_color, font_size_pt=9,
                 bold=False, letter_spacing_pt=6)

    # Textes bas du template Chanel — police/espacement originaux préservés
    _recolor_cover_texts(slide, text_color, text_color)
    _bring_cover_texts_to_front(slide)


def _cover_split(slide, W, H, primary: str, secondary: str,
                 logo_bytes: Optional[bytes], brand_name: str) -> None:
    """
    Style B — "split" : bandeau vertical couleur primaire à gauche (42 %),
    fond secondaire (ou blanc cassé) à droite. Logo centré sur le bandeau.
    """
    band_w    = int(W * 0.42)
    right_bg  = secondary if secondary != primary else "F8F6F2"
    on_dark_l = _is_dark(primary)
    on_dark_r = _is_dark(right_bg)
    text_l    = "FFFFFF" if on_dark_l else "111111"
    text_r    = "111111" if not on_dark_r else "FFFFFF"

    # Fond droite
    _set_slide_bg(slide, right_bg)

    # Bandeau gauche
    _add_rect(slide, 0, 0, band_w, H, primary)

    # Séparateur 1 px entre les deux zones
    _add_rect(slide, band_w, 0, int(W * 0.003), H, primary)

    # Logo centré dans le bandeau gauche — recolorisé pour contraster avec la bande
    if logo_bytes:
        logo_colored = _recolor_logo(logo_bytes, primary)
        _add_logo(slide, logo_colored,
                  cx=band_w // 2, cy=H // 2,
                  max_w=int(band_w * 0.68), max_h=int(H * 0.38))

    # Nom de marque à droite du bandeau, large et aéré
    _add_textbox(slide,
                 left=int(W * 0.47), top=int(H * 0.34),
                 width=int(W * 0.48), height=int(H * 0.12),
                 text=brand_name.upper(),
                 hex6_color=text_r, font_size_pt=18,
                 bold=True, letter_spacing_pt=8)

    # Textes bas du template Chanel — TextBox 5 (gauche/bande → text_l), TextBox 4 (droite → text_r)
    _recolor_cover_texts(slide, text_l, text_r)
    _bring_cover_texts_to_front(slide)

    # Petite ligne déco bas du bandeau
    _add_rect(slide, int(band_w * 0.12), int(H * 0.82),
              int(band_w * 0.76), int(H * 0.004), text_l)


def _cover_minimal(slide, W, H, primary: str, secondary: str,
                   logo_bytes: Optional[bytes], brand_name: str) -> None:
    """
    Style C — "minimal" : fond clair/secondaire, large barre primaire en bas (28 %).
    Logo centré dans la zone haute. Très sobre.
    """
    light_bg  = secondary if not _is_dark(secondary) else "F5F3EE"
    text_top  = "111111" if not _is_dark(light_bg) else "FFFFFF"

    # Couleur de la barre : primaire si visible, sinon secondaire (ex: Netflix noir → rouge)
    # Critère : luminance < 0.12 = trop sombre, on prend la secondaire si elle est plus lisible
    def _lum_hex(h6):
        r, g, b = int(h6[:2],16), int(h6[2:4],16), int(h6[4:],16)
        return (0.299*r + 0.587*g + 0.114*b) / 255
    bar_color = primary
    if _lum_hex(primary) < 0.12 and not _is_dark(secondary):
        bar_color = secondary  # ex: primaire=noir, secondaire=rouge → barre rouge
    elif _lum_hex(primary) < 0.12:
        # Les deux sont sombres → on prend primaire mais avec une hauteur réduite
        bar_color = primary
    bar_h     = int(H * 0.28)
    on_dark_b = _is_dark(bar_color)
    text_b    = "FFFFFF" if on_dark_b else "111111"

    # Fond clair
    _set_slide_bg(slide, light_bg)

    # Barre en bas (couleur intelligente)
    _add_rect(slide, 0, H - bar_h, W, bar_h, bar_color)

    # Logo centré dans la zone haute — couleurs originales préservées
    # Zone utile : dessous du brand name (H*0.13) jusqu'au haut de la barre (H-bar_h)
    # Centre optique = milieu de cette zone
    zone_top  = int(H * 0.13)
    zone_bot  = H - bar_h
    logo_cy   = (zone_top + zone_bot) // 2
    if logo_bytes:
        _add_logo(slide, logo_bytes,
                  cx=W // 2, cy=logo_cy,
                  max_w=int(W * 0.35), max_h=int((zone_bot - zone_top) * 0.55))

    # Nom de marque discret au-dessus du logo
    _add_textbox(slide,
                 left=int(W * 0.08), top=int(H * 0.06),
                 width=int(W * 0.84), height=int(H * 0.07),
                 text=brand_name.upper(),
                 hex6_color=text_top, font_size_pt=9,
                 bold=False, letter_spacing_pt=6)

    # Textes bas du template Chanel — dans la barre
    _recolor_cover_texts(slide, text_b, text_b)
    _bring_cover_texts_to_front(slide)


# ── Dispatcher principal ───────────────────────────────────────────────────

def build_cover(prs, req, logo_bytes: Optional[bytes]) -> None:
    """
    Génère la cover (slide 0) selon cover_style :
      "brand"     → Style A : fond plein couleur primaire, logo centré
      "split"     → Style B : bandeau vertical primaire + fond secondaire
      "minimal"   → Style C : fond clair, barre primaire en bas

    Dans tous les cas : pas d'image IA, uniquement logo + géométrie + couleurs.
    """
    slide = prs.slides[0]
    W, H  = prs.slide_width, prs.slide_height

    primary   = _hex(req.primary_color,   "1A1A1A")
    secondary = _hex(req.secondary_color or req.primary_color, "F5F3EE")
    brand_name = req.brand_name or ""
    style      = (req.cover_style or "brand").lower().strip()

    # Masquer les formes Chanel originales (on repart d'une page vierge)
    _hide_existing_cover_shapes(slide)

    # Dispatcher
    if style == "split":
        _cover_split(slide, W, H, primary, secondary, logo_bytes, brand_name)
    elif style == "minimal":
        _cover_minimal(slide, W, H, primary, secondary, logo_bytes, brand_name)
    else:  # "brand" (défaut)
        _cover_brand(slide, W, H, primary, secondary, logo_bytes, brand_name)


# ─── Cabine carrée : extraction + branding Pillow ─────────────────────────────

def _extract_shape_image(pptx_path: str, slide_idx: int, shape_name: str) -> Optional[bytes]:
    """Extracts the raw image bytes from a shape's blip in a PPTX file."""
    try:
        from pptx import Presentation as _Presentation
        _prs = _Presentation(pptx_path)
        _slide = _prs.slides[slide_idx]
        _shape = get_shape(_slide, shape_name)
        if _shape is None:
            return None
        blip = _shape.element.find('.//' + qn('a:blip'))
        if blip is None:
            return None
        rId = blip.get(f'{{{NS_R}}}embed')
        if not rId:
            return None
        return _slide.part.related_parts[rId].blob
    except Exception as e:
        print(f"_extract_shape_image({shape_name}) failed: {e}")
        return None


# Chemins vers les modèles neutres
MODELE_CARRE_PATH = os.path.join(os.path.dirname(__file__), "..", "..", "exemple", "modeleclassiquecarre.jpeg")
MODELE_ROND_PATH  = os.path.join(os.path.dirname(__file__), "..", "..", "exemple", "modeleclassiquerond.jpeg")


def _brand_cabine_carre(
    primary_hex: str,
    logo_bytes: Optional[bytes],
    slot_images: Optional[list] = None,   # list of up to 3 bytes objects
    brand_name: str = "",
) -> bytes:
    """
    Rebrande le modèle neutre carrée (modeleclassiquecarre.jpeg) :
    • Grand panneau DROIT  → couleur primaire (semi-transparent) + logo centré
    • 3 slots BLANCS gauche → photos brand (sans fond ajouté)
    • Texte "BOOTH" recouvert par le nom de marque (police système clean)
    Coordonnées calibrées sur l'image 1024×1022.
    """
    from PIL import Image as _Image, ImageDraw as _IDraw, ImageFont as _IFont
    import io as _io

    with open(MODELE_CARRE_PATH, "rb") as f:
        base_bytes = f.read()

    img = _Image.open(_io.BytesIO(base_bytes)).convert("RGBA")
    W, H = img.size  # 1024 × 1022

    hex_clean = primary_hex.lstrip("#")
    r, g, b = int(hex_clean[0:2], 16), int(hex_clean[2:4], 16), int(hex_clean[4:6], 16)

    # ── Panneau DROIT : fond couleur primaire semi-transparent ──────────────
    # Le panneau blanc droit est ~x: 51%-70%, y: 13%-83%
    px1, py1 = int(W * 0.512), int(H * 0.130)
    px2, py2 = int(W * 0.700), int(H * 0.835)
    overlay = _Image.new("RGBA", img.size, (0, 0, 0, 0))
    _IDraw.Draw(overlay).rectangle([px1, py1, px2, py2], fill=(r, g, b, 210))
    img = _Image.alpha_composite(img, overlay)

    # Logo centré sur le panneau droit
    if logo_bytes:
        try:
            logo = _Image.open(_io.BytesIO(logo_bytes)).convert("RGBA")
            pw, ph = px2 - px1, py2 - py1
            max_w, max_h = int(pw * 0.70), int(ph * 0.35)
            sc = min(max_w / logo.width, max_h / logo.height)
            logo = logo.resize((max(1, int(logo.width * sc)), max(1, int(logo.height * sc))), _Image.LANCZOS)
            img.paste(logo, (px1 + (pw - logo.width) // 2, py1 + (ph - logo.height) // 2), logo)
        except Exception as e:
            print(f"_brand_cabine_carre logo paste failed: {e}")

    # ── 3 slots GAUCHE : photos brand ───────────────────────────────────────
    # Coordonnées calibrées sur l'image 1024×1022 :
    #   Slot 1 (haut-gauche) :  x: 12.2%-21.0%, y: 21.7%-29.4%
    #   Slot 2 (haut-droite) :  x: 22.0%-29.4%, y: 21.7%-29.4%
    #   Slot 3 (bas, plus grand): x: 12.2%-21.0%, y: 30.5%-45.2%
    slots = [
        (int(W*0.122), int(H*0.217), int(W*0.210), int(H*0.294)),  # haut-gauche
        (int(W*0.220), int(H*0.217), int(W*0.294), int(H*0.294)),  # haut-droite
        (int(W*0.122), int(H*0.305), int(W*0.210), int(H*0.452)),  # bas (plus haut)
    ]
    if slot_images:
        for i, (sx1, sy1, sx2, sy2) in enumerate(slots):
            if i >= len(slot_images) or not slot_images[i]:
                continue
            try:
                simg = _Image.open(_io.BytesIO(slot_images[i])).convert("RGBA")
                sw, sh = sx2 - sx1, sy2 - sy1
                # Cover crop : on remplit le slot sans bande noire
                src_ratio = simg.width / simg.height
                slot_ratio = sw / sh
                if src_ratio > slot_ratio:
                    new_h = sh
                    new_w = int(sh * src_ratio)
                else:
                    new_w = sw
                    new_h = int(sw / src_ratio)
                simg = simg.resize((new_w, new_h), _Image.LANCZOS)
                # Centrer et cropper
                cx_off = (new_w - sw) // 2
                cy_off = (new_h - sh) // 2
                simg = simg.crop((cx_off, cy_off, cx_off + sw, cy_off + sh))
                img.paste(simg, (sx1, sy1), simg)
            except Exception as e:
                print(f"_brand_cabine_carre slot {i} failed: {e}")

    # ── Nom de marque par-dessus "BOOTH" ─────────────────────────────────────
    if brand_name:
        try:
            draw = _IDraw.Draw(img)
            # "BOOTH" est ~x: 13%-27%, y: 16%-21%
            tx, ty = int(W * 0.130), int(H * 0.160)
            tbox_w, tbox_h = int(W * 0.150), int(H * 0.055)
            # Peindre par-dessus le texte Booth existant (même couleur bois ~sombre)
            # On overlay avec la couleur primaire
            text_bg_overlay = _Image.new("RGBA", img.size, (0, 0, 0, 0))
            _IDraw.Draw(text_bg_overlay).rectangle(
                [tx - 4, ty - 2, tx + tbox_w + 4, ty + tbox_h + 2],
                fill=(r, g, b, 230)
            )
            img = _Image.alpha_composite(img, text_bg_overlay)
            # Texte blanc
            try:
                font_size = max(14, int(tbox_h * 0.65))
                font = _IFont.truetype("/System/Library/Fonts/Helvetica.ttc", font_size)
            except Exception:
                font = _IFont.load_default()
            draw = _IDraw.Draw(img)
            brand_upper = brand_name.upper()
            try:
                bbox = draw.textbbox((0, 0), brand_upper, font=font)
                tw = bbox[2] - bbox[0]
            except Exception:
                tw = len(brand_upper) * font_size // 2
            draw.text((tx + (tbox_w - tw) // 2, ty + (tbox_h - font_size) // 2),
                      brand_upper, fill=(255, 255, 255, 255), font=font)
        except Exception as e:
            print(f"_brand_cabine_carre brand name overlay failed: {e}")

    out = _io.BytesIO()
    img.convert("RGB").save(out, format="PNG")
    return out.getvalue()


def _brand_cabine_ronde(
    primary_hex: str,
    logo_bytes: Optional[bytes],
    brand_name: str = "",
) -> bytes:
    """
    Rebrande le modèle neutre rond (modeleclassiquerond.jpeg) :
    • Panneau principal → couleur primaire semi-transparent + logo centré
    Coordonnées calibrées sur l'image 819×1024.
    """
    from PIL import Image as _Image, ImageDraw as _IDraw, ImageFont as _IFont
    import io as _io

    with open(MODELE_ROND_PATH, "rb") as f:
        base_bytes = f.read()

    img = _Image.open(_io.BytesIO(base_bytes)).convert("RGBA")
    W, H = img.size  # 819 × 1024

    hex_clean = primary_hex.lstrip("#")
    r, g, b = int(hex_clean[0:2], 16), int(hex_clean[2:4], 16), int(hex_clean[4:6], 16)

    # Panneau principal de la cabine ronde : ~x: 10%-85%, y: 8%-88%
    px1, py1 = int(W * 0.10), int(H * 0.08)
    px2, py2 = int(W * 0.85), int(H * 0.88)
    overlay = _Image.new("RGBA", img.size, (0, 0, 0, 0))
    _IDraw.Draw(overlay).rectangle([px1, py1, px2, py2], fill=(r, g, b, 180))
    img = _Image.alpha_composite(img, overlay)

    # Logo centré sur le panneau
    if logo_bytes:
        try:
            logo = _Image.open(_io.BytesIO(logo_bytes)).convert("RGBA")
            pw, ph = px2 - px1, py2 - py1
            max_w, max_h = int(pw * 0.60), int(ph * 0.25)
            sc = min(max_w / logo.width, max_h / logo.height)
            logo = logo.resize((max(1, int(logo.width * sc)), max(1, int(logo.height * sc))), _Image.LANCZOS)
            img.paste(logo, (px1 + (pw - logo.width) // 2, py1 + (ph - logo.height) // 2), logo)
        except Exception as e:
            print(f"_brand_cabine_ronde logo paste failed: {e}")

    out = _io.BytesIO()
    img.convert("RGB").save(out, format="PNG")
    return out.getvalue()


@router.post("/generate-presentation")
def generate_presentation(req: GenerateRequest):
    # Base = Chanel (layout correct : 4 colonnes cabines, 3 colonnes kiosk, 2 goodies)
    # Fallback = template neutre si Chanel introuvable
    base_path = CHANEL_PATH if os.path.exists(CHANEL_PATH) else TEMPLATE_PATH
    if not os.path.exists(base_path):
        raise HTTPException(status_code=404, detail=f"Template not found: {base_path}")

    prs = Presentation(base_path)

    # Toutes les slides Chanel sont conservées (aucune suppression)

    # ── Logo de la marque (utilisé cover + strips photo) ─────────────────────
    logo_bytes: Optional[bytes] = None
    if req.logo_url:
        try:
            logo_bytes = get_image_bytes(req.logo_url)
        except Exception as e:
            print(f"Warning: logo fetch failed: {e}")
    # Icône/symbole séparé (ex: le «N» de Netflix pour Freeform 7)
    # Si absent, on réutilise le logo principal
    logo_icon_bytes: Optional[bytes] = None
    if req.logo_icon_url:
        try:
            logo_icon_bytes = get_image_bytes(req.logo_icon_url)
        except Exception as e:
            print(f"Warning: logo_icon fetch failed: {e}")
    if logo_icon_bytes is None:
        logo_icon_bytes = logo_bytes  # fallback : même logo que Freeform 5
    # ── Zones brand à remplacer via replace_blip ──────────────────────────────
    # Slide 0  : Freeform 3 (top-level) = logo marque cover
    # Slide 3  (p4 "Photobooth classique") : deux PETITES images de la bande
    #           → Freeform 5 (enfant Group 4) et Freeform 7 (enfant Group 6)
    #           → NE PAS toucher Freeform 3 / Group 2 (la grande photo, figée)
    # Slide 4  (p5 "Cabines photos") : Freeform 25 = cabine top, Freeform 24 = bas
    # Slide 5  (p6 "The Kiosk") : Freeform 8 = photo éditoriale cabine
    # Goodies  (slide 9 après suppression = p11 "Nos goodies") :
    #           → inject_picture_at_shape sur "Group 2" et "Group 4"
    #             (les enfants Freeform 3/5 sont minuscules – l'image est dans le groupe)
    # Cover (slide 0) est géré séparément par build_cover() — pas de blip ici.
    blip_zones = [
        # Slide 3 Freeform 5 & 7 : PAS d'injection IA — vraies photos Chanel conservées.
        # Le branding (nom marque + logo) est appliqué via _rebrand_photo_strips().
        (5, "Freeform 8", req.kiosk_url),  # p6 kiosk
    ]

    # ── Téléchargements — images kiosk uniquement ────────────────────────────
    import time
    fetched: dict[str, bytes] = {}
    for idx, name, url in blip_zones:
        if not url:
            continue
        key = f"{idx}_{name}"
        try:
            fetched[key] = get_image_bytes(url)
        except Exception as e:
            print(f"Warning: failed to fetch '{key}': {e}")

    # ── Appliquer ─────────────────────────────────────────────────────────────
    for idx, name, url in blip_zones:
        if not url:
            continue
        key = f"{idx}_{name}"
        img_bytes = fetched.get(key)
        if not img_bytes:
            continue
        slide = prs.slides[idx]
        shape = get_shape(slide, name)
        if shape is None:
            print(f"Warning: shape '{name}' not found on slide {idx}")
            continue
        # Use replace_blip_contained for kiosk to avoid stretching
        if name == "Freeform 8":
            if not replace_blip_contained(slide.part, shape, img_bytes):
                print(f"Warning: replace_blip_contained failed for '{name}' on slide {idx}")
        else:
            if not replace_blip(slide.part, shape, img_bytes):
                print(f"Warning: replace_blip failed for '{name}' on slide {idx}")

    # ── Slots images cabine carrée : téléchargement ──────────────────────────
    slot_images: list = []
    if req.cabine_slot_urls:
        for i, slot_url in enumerate(req.cabine_slot_urls[:3]):
            if not slot_url:
                slot_images.append(None)
                continue
            try:
                slot_images.append(get_image_bytes(slot_url))
            except Exception as e:
                print(f"Warning: cabine slot {i} fetch failed: {e}")
                slot_images.append(None)

    primary_h = _hex(req.primary_color, "1A1A1A")
    slide4 = prs.slides[4]

    # ── Cabine CARRÉE (Freeform 24) ──────────────────────────────────────────
    # Priority: AI image from Gemini → fallback to Pillow local model
    try:
        if req.cabine_carree_url:
            branded_f24 = get_image_bytes(req.cabine_carree_url)
        else:
            branded_f24 = _brand_cabine_carre(
                primary_hex  = primary_h,
                logo_bytes   = logo_bytes,
                slot_images  = slot_images or None,
                brand_name   = req.brand_name or "",
            )
        shape_f24 = get_shape(slide4, "Freeform 24")
        if shape_f24 is not None:
            if not replace_blip(slide4.part, shape_f24, branded_f24):
                print("Warning: replace_blip failed for Freeform 24")
        else:
            print("Warning: Freeform 24 not found on slide 4")
    except Exception as e:
        print(f"Warning: cabine carrée (Freeform 24) failed: {e}")

    # ── Cabine RONDE (Freeform 25) ───────────────────────────────────────────
    # Priority: AI image from Gemini → fallback to Pillow local model
    try:
        if req.cabine_ronde_url:
            branded_f25 = get_image_bytes(req.cabine_ronde_url)
        else:
            branded_f25 = _brand_cabine_ronde(
                primary_hex = primary_h,
                logo_bytes  = logo_bytes,
                brand_name  = req.brand_name or "",
            )
        shape_f25 = get_shape(slide4, "Freeform 25")
        if shape_f25 is not None:
            if not replace_blip(slide4.part, shape_f25, branded_f25):
                print("Warning: replace_blip failed for Freeform 25")
        else:
            print("Warning: Freeform 25 not found on slide 4")
    except Exception as e:
        print(f"Warning: cabine ronde (Freeform 25) failed: {e}")

    # ── Bandes photo slide 3 : rebranding CHANEL → marque (nom + logo) ──────────
    try:
        _rebrand_photo_strips(
            prs.slides[3],
            primary    = _hex(req.primary_color,                         "1A1A1A"),
            secondary  = _hex(req.secondary_color or req.primary_color,  "F5F3EE"),
            brand_name = req.brand_name or "",
            logo_bytes = logo_bytes,       # wordmark → Freeform 5 (fond primaire)
            logo_icon_bytes = logo_icon_bytes,  # icône → Freeform 7 (fond secondaire)
        )
    except Exception as e:
        print(f"Warning: _rebrand_photo_strips failed: {e}")
    # ── Bandes photo slide 3 : contour de marque sur Freeform 5 & 7 ────────────
    try:
        _style_photo_strips(
            prs.slides[3],
            primary   = _hex(req.primary_color,              "1A1A1A"),
            secondary = _hex(req.secondary_color or req.primary_color, "F5F3EE"),
            strip_style = req.strip_style or "none",
        )
    except Exception as e:
        print(f"Warning: _style_photo_strips failed: {e}")

    # ── Slide tarifs (index 12) : fond couleur primaire + logo watermark ──────
    try:
        _rebrand_tarifs_slide(
            prs,
            slide_idx = 12,
            primary_hex = _hex(req.primary_color, "1A1A1A"),
            logo_bytes  = logo_bytes,
        )
    except Exception as e:
        print(f'Warning: _rebrand_tarifs_slide failed: {e}')

    # ── Goodies : inject_picture_at_shape sur le groupe (overlay fiable) ─────
    # Chanel template : slide index 10 = "Nos goodies" (aucune suppression en amont).
    GOODIES_SLIDE = 10
    goodies_inject = [
        ("Group 2", req.goodies_top_url),
        ("Group 4", req.goodies_bottom_url),
    ]
    for group_name, url in goodies_inject:
        if not url:
            continue
        try:
            img_bytes_g = get_image_bytes(url)
        except Exception as e:
            print(f"Warning: goodies fetch failed for '{group_name}': {e}")
            continue
        try:
            inject_picture_at_shape(prs.slides[GOODIES_SLIDE], group_name, img_bytes_g)
        except Exception as e:
            print(f"Warning: goodies inject failed for '{group_name}': {e}")

    # ── Cover : logo centré + fond uni (3 styles, pas d'image IA) ───────────
    build_cover(prs, req, logo_bytes)

    # ── Remplacement du texte de marque ───────────────────────────────────────
    brand_title = req.brand_name.title()
    brand_upper = req.brand_name.upper()
    website = req.website or ""

    for slide in prs.slides:
        replace_text(slide, "CHANEL", brand_upper)
        replace_text(slide, "Chanel", brand_title)
        replace_text(slide, "chanel", req.brand_name.lower())
        if website:
            replace_text(slide, "chanel.com", website)

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)

    fname = req.brand_name.replace(" ", "_") + "_x_Booth.pptx"
    return StreamingResponse(
        out,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{fname}"'},
    )


# ── Renderer Pillow : slide → PNG ──────────────────────────────────────────

def _render_slide_png(slide, prs, thumb_w: int = 1280) -> bytes:
    """
    Rendu basique d'un slide en PNG via Pillow.
    Gère : fond, rectangles colorés, images (blip), groupes.
    Suffisant pour un aperçu de développement.
    """
    from PIL import Image as _PIL, ImageDraw as _IDraw, ImageFont as _IFont
    from pptx.enum.dml import MSO_THEME_COLOR
    from pptx.dml.color import RGBColor

    W_emu = prs.slide_width
    H_emu = prs.slide_height
    scale  = thumb_w / W_emu
    thumb_h = int(H_emu * scale)

    # --- fond ---
    bg_color = (255, 255, 255)
    try:
        bg_fill = slide.background.fill
        if bg_fill.type is not None:
            c = bg_fill.fore_color.rgb
            bg_color = (c.r, c.g, c.b)
    except Exception:
        pass
    canvas = _PIL.new("RGB", (thumb_w, thumb_h), bg_color)

    def _emu_rect(shape):
        x = int(shape.left  * scale)
        y = int(shape.top   * scale)
        w = int(shape.width * scale)
        h = int(shape.height * scale)
        return x, y, w, h

    def _try_get_blip_bytes(shape):
        sp = shape._element
        blip = sp.find('.//' + qn('a:blip'))
        if blip is None:
            return None
        rId = blip.get(f'{{{NS_R}}}embed')
        if not rId:
            return None
        try:
            for rel in slide.part.rels.values():
                if rel.rId == rId and rel.reltype.endswith('/image'):
                    return rel._target.blob
        except Exception:
            pass
        return None

    def _render_shape(shape):
        try:
            x, y, w, h = _emu_rect(shape)
            if w <= 0 or h <= 0:
                return
        except Exception:
            return

        # Groupe → recurse
        if shape.shape_type == 6:
            for child in shape.shapes:
                _render_shape(child)
            return

        draw = _IDraw.Draw(canvas)

        # Image (blip)
        blob = _try_get_blip_bytes(shape)
        if blob:
            try:
                img = _PIL.open(io.BytesIO(blob)).convert("RGBA")
                img = img.resize((max(1, w), max(1, h)), _PIL.LANCZOS)
                canvas.paste(img, (x, y), img)
                return
            except Exception:
                pass

        # Solid fill
        try:
            fill = shape.fill
            if fill.type is not None:
                try:
                    c = fill.fore_color.rgb
                    draw.rectangle([x, y, x + w, y + h], fill=(c.r, c.g, c.b))
                    return
                except Exception:
                    pass
        except Exception:
            pass

        # Text box avec fill transparent → dessine le texte
        try:
            if shape.has_text_frame:
                tf = shape.text_frame
                txt = tf.text.strip()
                if txt:
                    font_size = max(10, int(12 * scale * 914400 / 12700))  # approx
                    try:
                        font = _IFont.truetype("/System/Library/Fonts/Helvetica.ttc", size=max(8, int(h * 0.4)))
                    except Exception:
                        font = _IFont.load_default()
                    draw.text((x + 4, y + 2), txt, fill=(80, 80, 80), font=font)
        except Exception:
            pass

    for shape in slide.shapes:
        _render_shape(shape)

    out = io.BytesIO()
    canvas.save(out, format="PNG", optimize=True)
    return out.getvalue()


@router.post("/preview-slides")
def preview_slides(req: GenerateRequest):
    """
    Génère le PPTX (cover + slide 3 rebranding, SANS images IA)
    et renvoie les thumbnails PNG en base64 — un par slide.
    Endpoint local dev uniquement, pas d'appel Imagen.
    """
    import json as _json

    prs = Presentation(CHANEL_PATH)

    # Toutes les slides Chanel sont conservées — aucune suppression

    # Logo
    logo_bytes: Optional[bytes] = None
    if req.logo_url:
        try:
            logo_bytes = get_image_bytes(req.logo_url)
        except Exception as e:
            print(f"Preview: logo fetch failed: {e}")

    logo_icon_bytes: Optional[bytes] = None
    if req.logo_icon_url:
        try:
            logo_icon_bytes = get_image_bytes(req.logo_icon_url)
        except Exception as e:
            print(f"Preview: logo_icon fetch failed: {e}")
    if logo_icon_bytes is None:
        logo_icon_bytes = logo_bytes

    # Slide 3 : rebranding strips
    try:
        _rebrand_photo_strips(
            prs.slides[3],
            primary    = _hex(req.primary_color,                        "1A1A1A"),
            secondary  = _hex(req.secondary_color or req.primary_color, "F5F3EE"),
            brand_name = req.brand_name or "",
            logo_bytes = logo_bytes,
            logo_icon_bytes = logo_icon_bytes,
        )
    except Exception as e:
        print(f"Preview: _rebrand_photo_strips failed: {e}")

    try:
        _style_photo_strips(
            prs.slides[3],
            primary     = _hex(req.primary_color,                         "1A1A1A"),
            secondary   = _hex(req.secondary_color or req.primary_color,  "F5F3EE"),
            strip_style = req.strip_style or "none",
        )
    except Exception as e:
        print(f"Preview: _style_photo_strips failed: {e}")

    # Slide tarifs : fond couleur primaire + logo watermark
    try:
        _rebrand_tarifs_slide(
            prs,
            slide_idx = 12,
            primary_hex = _hex(req.primary_color, "1A1A1A"),
            logo_bytes  = logo_bytes,
        )
    except Exception as e:
        print(f'Preview: _rebrand_tarifs_slide failed: {e}')

    # Cover
    build_cover(prs, req, logo_bytes)

    # Remplacement texte
    brand_title = req.brand_name.title()
    brand_upper = req.brand_name.upper()
    for slide in prs.slides:
        replace_text(slide, "CHANEL", brand_upper)
        replace_text(slide, "Chanel", brand_title)
        replace_text(slide, "chanel", req.brand_name.lower())

    # Render chaque slide en PNG
    slides_b64 = []
    for slide in prs.slides:
        try:
            png = _render_slide_png(slide, prs, thumb_w=1280)
            slides_b64.append("data:image/png;base64," + base64.b64encode(png).decode())
        except Exception as e:
            print(f"Preview: render failed for slide: {e}")
            slides_b64.append(None)

    from fastapi.responses import JSONResponse
    return JSONResponse({"slides": slides_b64})


@router.get("/health")
def health():
    return {
        "status": "ok",
        "template_exists": os.path.exists(TEMPLATE_PATH),
    }

